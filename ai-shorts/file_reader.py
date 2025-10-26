import os
from pypdf import PdfReader
from pptx import Presentation
from docx import Document


def extract_pdf(path):
    reader = PdfReader(path)
    extracted_text = []

    for i, page in enumerate(reader.pages, start=1):
        text = page.extract_text()
        if text:
            extracted_text.append(f"\n--- Page {i} ---\n{text.strip()}")
    return "\n".join(extracted_text)


def extract_pptx(path):
    presentation = Presentation(path)
    extracted_text = []

    for slide_number, slide in enumerate(presentation.slides, start=1):
        slide_text = [f"\n--- Slide {slide_number} ---"]
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text.strip())
        extracted_text.append("\n".join(slide_text))
    return "\n".join(extracted_text)


def extract_docx(path):
    doc = Document(path)
    paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    return "\n".join(paragraphs)


def extract_txt(path):
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()


def extract_text(path: str):
    ext = os.path.splitext(path)[1].lower()

    handlers = {
        ".pdf": extract_pdf,
        ".pptx": extract_pptx,
        ".docx": extract_docx,
        ".txt": extract_txt,
        ".html": extract_txt,
        ".htm": extract_txt,
        ".md": extract_txt,
    }

    if ext not in handlers:
        raise ValueError(f"Unsupported file type: {ext}")

    return handlers[ext](path)
